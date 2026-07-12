# /// script
# requires-python = ">=3.12"
# dependencies = ["python-pptx>=1.0.2", "pillow>=11.0.0"]
# ///

from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
MEDIA = ROOT / "presentation" / "media"
OUTPUT = ROOT / "presentation" / "output"
OUTPUT.mkdir(parents=True, exist_ok=True)

# 16:9 canvas
SW, SH = 13.333, 7.5

# Cerno design tokens
INK = "182320"
MUTED = "64726D"
PORCELAIN = "F3F6F5"
PAPER = "FCFDFC"
GREEN = "1C665A"
GREEN_DARK = "102A24"
SEA = "DCEBE7"
SEA_LIGHT = "EEF6F3"
AMBER = "D79A42"
OXIDE = "A45C52"
LINE = "D5DEDA"
WHITE = "FFFFFF"
SOFT = "E8EFED"

FONT_UI = "Avenir Next"
FONT_SERIF = "Georgia"
FONT_MONO = "Menlo"


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


def crop_image(src: Path, dst: Path, box):
    with Image.open(src) as image:
        image.crop(box).save(dst, quality=95)


# Purpose-built crops from real application screenshots.
briefing_src = ROOT / "assets" / "cerno-live-briefing.png"
crop_image(briefing_src, MEDIA / "briefing-demo-crop.png", (215, 75, 1590, 1010))
crop_image(briefing_src, MEDIA / "evidence-inspector.png", (980, 95, 1580, 795))


prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
blank = prs.slide_layouts[6]


def fill_background(slide, color=PORCELAIN):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb(color)


def add_rect(slide, x, y, w, h, fill=PAPER, line=None, radius=True, transparency=0):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.fill.transparency = transparency
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(0.8)
    else:
        shape.line.fill.background()
    return shape


def add_line(slide, x1, y1, x2, y2, color=LINE, width=1.0):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    return line


def add_text(slide, text, x, y, w, h, *, size=18, color=INK, font=FONT_UI,
             bold=False, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
             margin=0, line_spacing=1.0, letter_spacing=None, italic=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.line_spacing = line_spacing
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    # PowerPoint represents explicit line breaks as separate runs; style every run.
    for run in p.runs:
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = rgb(color)
    if letter_spacing is not None:
        # Not exposed by python-pptx; kept as semantic input for consistent call sites.
        pass
    return box


def add_label(slide, text, x, y, w, color=GREEN):
    return add_text(slide, text.upper(), x, y, w, 0.25, size=9, color=color,
                    font=FONT_MONO, bold=True, valign=MSO_ANCHOR.MIDDLE)


def add_logo(slide, x=0.62, y=0.42, inverse=False):
    c = "9CE1D3" if inverse else GREEN
    tc = WHITE if inverse else INK
    add_line(slide, x, y + 0.02, x + 0.28, y + 0.02, c, 1.6)
    add_line(slide, x + 0.045, y + 0.10, x + 0.235, y + 0.10, c, 1.6)
    add_line(slide, x + 0.095, y + 0.18, x + 0.185, y + 0.18, c, 1.6)
    add_text(slide, "CERNO", x + 0.40, y - 0.03, 1.15, 0.28, size=13, color=tc, bold=True,
             valign=MSO_ANCHOR.MIDDLE)


def add_footer(slide, number, dark=False, appendix=False):
    color = "AFC0BA" if dark else MUTED
    left = "CERNO · AI AS AGENCY"
    if appendix:
        left += " · APPENDIX"
    add_text(slide, left, 0.65, 7.17, 3.4, 0.18, size=7.5, color=color, font=FONT_MONO)
    add_text(slide, f"{number:02d}", 12.22, 7.15, 0.45, 0.18, size=8, color=color,
             font=FONT_MONO, align=PP_ALIGN.RIGHT)


def add_picture_contain(slide, path, x, y, w, h, fill=PAPER, border=LINE, pad=0.08):
    add_rect(slide, x, y, w, h, fill=fill, line=border, radius=True)
    with Image.open(path) as image:
        iw, ih = image.size
    scale = min((w - 2 * pad) / iw, (h - 2 * pad) / ih)
    pw, ph = iw * scale, ih * scale
    pic = slide.shapes.add_picture(str(path), Inches(x + (w - pw) / 2), Inches(y + (h - ph) / 2),
                                   width=Inches(pw), height=Inches(ph))
    return pic


def metric_card(slide, x, y, w, value, label, detail=None, accent=GREEN):
    add_rect(slide, x, y, w, 1.12, fill=PAPER, line=LINE, radius=True)
    add_rect(slide, x, y, 0.07, 1.12, fill=accent, radius=False)
    add_text(slide, value, x + 0.22, y + 0.18, w - 0.35, 0.43, size=24, color=INK, bold=True,
             valign=MSO_ANCHOR.MIDDLE)
    add_label(slide, label, x + 0.22, y + 0.68, w - 0.35, color=MUTED)
    if detail:
        add_text(slide, detail, x + 0.22, y + 0.90, w - 0.35, 0.16, size=7, color=MUTED)


# Slide 1 — Cover
slide = prs.slides.add_slide(blank)
fill_background(slide, GREEN_DARK)
add_logo(slide, inverse=True)
add_text(slide, "GROWTHX × HERMES BUILDATHON", 9.25, 0.47, 3.4, 0.2, size=8.5,
         color="AFC0BA", font=FONT_MONO, bold=True, align=PP_ALIGN.RIGHT)
add_text(slide, "Your personal research desk.\nStaffed by agents.", 0.72, 1.02, 11.5, 1.4,
         size=36, color=WHITE, font=FONT_SERIF, bold=True, line_spacing=0.92)
add_text(slide, "One mission in. One finite, cited briefing out.", 0.76, 2.48, 8.3, 0.36,
         size=14, color="C8D8D4")
slide.shapes.add_picture(str(MEDIA / "cerno-banner.png"), Inches(0), Inches(2.88),
                         width=Inches(SW), height=Inches(4.67))
# Banner has rounded internal edges; cover the tiny overflow at bottom.
add_rect(slide, 0, 7.42, SW, 0.08, fill=GREEN_DARK, radius=False)


# Slide 2 — Problem
slide = prs.slides.add_slide(blank)
fill_background(slide)
add_logo(slide)
add_label(slide, "The problem", 0.72, 0.98, 2.0)
add_text(slide, "The internet does not have an\ninformation shortage.", 0.72, 1.30, 11.8, 1.02,
         size=31, font=FONT_SERIF, bold=True, line_spacing=0.95)
add_text(slide, "It has a judgment problem.", 0.72, 2.28, 11.8, 0.52,
         size=24, color=GREEN, font=FONT_SERIF, bold=True)
add_line(slide, 0.72, 3.02, 12.60, 3.02, LINE, 1.0)
problem_items = [
    ("01", "FEEDS", "Optimize engagement —\nnot the work in front of you."),
    ("02", "READ-LATER", "Collect more material —\nand leave the reading to you."),
    ("03", "SUMMARIES", "Make everything shorter —\nwithout deciding if it matters."),
]
for i, (num, title, body) in enumerate(problem_items):
    x = 0.72 + i * 4.02
    add_text(slide, num, x, 3.40, 0.42, 0.28, size=10, color=GREEN, font=FONT_MONO, bold=True)
    add_text(slide, title, x + 0.50, 3.38, 2.8, 0.28, size=11, color=INK, bold=True)
    add_text(slide, body, x + 0.50, 3.83, 3.15, 0.76, size=15, color=MUTED, line_spacing=1.12)
add_rect(slide, 0.72, 5.15, 11.88, 1.16, fill=GREEN_DARK, radius=True)
add_text(slide, "The missing function", 1.04, 5.43, 2.2, 0.22, size=9, color="9CE1D3", font=FONT_MONO, bold=True)
add_text(slide, "Decide what deserves attention now — for this person.", 3.20, 5.32, 8.85, 0.52,
         size=21, color=WHITE, font=FONT_SERIF, bold=True, valign=MSO_ANCHOR.MIDDLE)
add_footer(slide, 2)


# Slide 3 — Product and demo transition
slide = prs.slides.add_slide(blank)
fill_background(slide)
add_logo(slide)
add_label(slide, "The product", 0.72, 0.93, 2.2)
add_text(slide, "One mission in. One finite briefing out.", 0.72, 1.25, 11.8, 0.60,
         size=28, font=FONT_SERIF, bold=True)
steps = [
    ("01", "FOCUS THREAD", "Declare the current decision, context, and boundaries."),
    ("02", "AGENT TEAM", "A Hermes Director delegates text, video, and personal judgment."),
    ("03", "EXACT EVIDENCE", "Passages or VideoDB transcript moments must exact-match before publication."),
    ("04", "VISIBLE TASTE", "Correct the reasoning; approve a readable TasteDoc change."),
]
for i, (num, title, body) in enumerate(steps):
    y = 2.05 + i * 1.03
    add_text(slide, num, 0.74, y + 0.03, 0.44, 0.24, size=10, color=GREEN, font=FONT_MONO, bold=True)
    add_text(slide, title, 1.28, y, 2.95, 0.24, size=10, color=INK, bold=True)
    add_text(slide, body, 1.28, y + 0.31, 3.75, 0.53, size=11.5, color=MUTED, line_spacing=1.08)
    if i < 3:
        add_line(slide, 0.95, y + 0.35, 0.95, y + 0.96, SEA, 1.6)
add_picture_contain(slide, MEDIA / "briefing-demo-crop.png", 5.28, 1.94, 7.30, 4.63,
                    fill="E5ECEA", border=LINE, pad=0.08)
add_rect(slide, 0.72, 6.38, 4.03, 0.42, fill=SEA, radius=True)
add_text(slide, "LIVE DEMO  →  assign · inspect · correct", 0.92, 6.47, 3.65, 0.18,
         size=8.5, color=GREEN, font=FONT_MONO, bold=True, valign=MSO_ANCHOR.MIDDLE)
add_footer(slide, 3)


# Slide 4 — Agency architecture
slide = prs.slides.add_slide(blank)
fill_background(slide)
add_logo(slide)
add_label(slide, "AI as Agency", 0.72, 0.93, 2.0)
add_text(slide, "An agent team replaces the personal research desk.", 0.72, 1.22, 11.8, 0.62,
         size=28, font=FONT_SERIF, bold=True)
# Input
add_rect(slide, 0.72, 2.23, 2.08, 2.88, fill=PAPER, line=LINE, radius=True)
add_label(slide, "Assignment", 0.96, 2.49, 1.2)
add_text(slide, "Focus Thread", 0.96, 2.82, 1.45, 0.28, size=16, bold=True)
add_text(slide, "Current mission\nKnown context\nDesired outcome\nFreshness\nTasteDoc snapshot", 0.96, 3.28, 1.55, 1.26,
         size=11.5, color=MUTED, line_spacing=1.34)
# Connect input to director
add_line(slide, 2.80, 3.64, 3.34, 3.64, GREEN, 1.8)
add_rect(slide, 3.34, 2.54, 2.18, 2.18, fill=GREEN_DARK, radius=True)
add_label(slide, "Manager agent", 3.66, 2.85, 1.50, color="9CE1D3")
add_text(slide, "Research\nDirector", 3.66, 3.20, 1.50, 0.78, size=22, color=WHITE,
         font=FONT_SERIF, bold=True, line_spacing=0.92)
add_text(slide, "plans · delegates · reviews", 3.66, 4.17, 1.60, 0.20, size=7.5,
         color="AFC0BA", font=FONT_MONO)
# Specialist fan-out
add_line(slide, 5.52, 3.63, 6.08, 3.63, GREEN, 1.8)
add_line(slide, 6.08, 2.64, 6.08, 4.65, GREEN, 1.2)
specialists = [
    (2.20, "EVIDENCE\nANALYST", "primary sources"),
    (3.25, "VIDEO\nANALYST", "timestamped moments"),
    (4.30, "PERSONAL\nEDITOR", "novelty + taste fit"),
]
for y, role, detail in specialists:
    add_line(slide, 6.08, y + 0.39, 6.42, y + 0.39, GREEN, 1.2)
    add_rect(slide, 6.42, y, 2.28, 0.78, fill=PAPER, line=LINE, radius=True)
    add_text(slide, role, 6.64, y + 0.12, 1.02, 0.42, size=9.5, color=INK, bold=True, line_spacing=0.9)
    add_text(slide, detail, 7.72, y + 0.20, 0.78, 0.32, size=7.5, color=MUTED)
# Output checks
add_line(slide, 8.70, 3.63, 9.22, 3.63, GREEN, 1.8)
add_rect(slide, 9.22, 2.23, 3.38, 2.88, fill=SEA_LIGHT, line=SEA, radius=True)
add_label(slide, "Deterministic gate", 9.50, 2.49, 1.9)
add_text(slide, "Validate before publish", 9.50, 2.82, 2.65, 0.30, size=16, bold=True)
checks = ["exact quote", "source or timestamp locator", "content hash", "playable video moment"]
for i, check in enumerate(checks):
    cy = 3.30 + i * 0.38
    add_rect(slide, 9.51, cy, 0.18, 0.18, fill=GREEN, radius=True)
    add_text(slide, "✓", 9.51, cy - 0.01, 0.18, 0.18, size=7, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, check, 9.82, cy - 0.01, 1.75, 0.20, size=10.5, color=MUTED)
add_text(slide, "FINITE BRIEFING", 9.50, 4.79, 2.50, 0.20, size=8, color=GREEN, font=FONT_MONO, bold=True)
# Technology boundary
add_rect(slide, 0.72, 5.63, 11.88, 0.78, fill=PAPER, line=LINE, radius=True)
techs = [
    ("CONVEX", "canonical memory"),
    ("LINKUP", "live discovery + fetch"),
    ("VIDEODB", "indexed exact moments"),
    ("HERMES", "Director + delegation"),
]
for i, (name, detail) in enumerate(techs):
    x = 0.98 + i * 2.92
    add_text(slide, name, x, 5.87, 0.78, 0.20, size=7.8, color=GREEN, font=FONT_MONO, bold=True)
    add_text(slide, detail, x + 0.82, 5.84, 1.98, 0.24, size=9.2, color=MUTED)
add_footer(slide, 4)


# Slide 5 — Trust / evidence spine
slide = prs.slides.add_slide(blank)
fill_background(slide)
add_logo(slide)
add_label(slide, "The trust loop", 0.72, 0.93, 2.2)
add_text(slide, "Trust is a product surface — not a disclaimer.", 0.72, 1.22, 11.8, 0.62,
         size=28, font=FONT_SERIF, bold=True)
add_text(slide, "Every finding exposes the full path from claim to personal judgment.", 0.74, 1.90, 6.0, 0.32,
         size=12.5, color=MUTED)
# Evidence spine
spine_x = 1.04
add_line(slide, spine_x, 2.66, spine_x, 5.93, GREEN, 2.0)
spine_items = [
    (2.52, "01", "CLAIM", "Atomic, decision-relevant finding"),
    (3.26, "02", "EXACT PASSAGE", "Source substring or timestamped VideoDB transcript"),
    (4.00, "03", "PROVENANCE", "Source/timestamp locator + SHA-256 content hash"),
    (4.74, "04", "PERSONAL INDEX", "Novelty and redundancy against prior claims"),
    (5.48, "05", "TASTEDOC RULE", "Visible rule used to make the judgment"),
]
for y, num, title, body in spine_items:
    add_rect(slide, spine_x - 0.13, y + 0.08, 0.26, 0.26, fill=GREEN, radius=True)
    add_text(slide, num, 1.36, y, 0.35, 0.22, size=8, color=GREEN, font=FONT_MONO, bold=True)
    add_text(slide, title, 1.80, y, 1.72, 0.22, size=9.5, color=INK, bold=True)
    add_text(slide, body, 1.80, y + 0.28, 3.85, 0.30, size=10.5, color=MUTED)
# Actual UI crop
add_picture_contain(slide, MEDIA / "evidence-inspector.png", 6.42, 1.92, 5.10, 4.86,
                    fill="E4ECE9", border=LINE, pad=0.07)
# Feedback loop badge
add_rect(slide, 9.72, 5.94, 2.65, 0.66, fill=GREEN_DARK, radius=True)
add_text(slide, "CORRECTION → REVIEWED DIFF", 9.90, 6.08, 2.30, 0.18, size=7.5,
         color="9CE1D3", font=FONT_MONO, bold=True, align=PP_ALIGN.CENTER)
add_footer(slide, 5)


# Slide 6 — Proof
slide = prs.slides.add_slide(blank)
fill_background(slide)
add_logo(slide)
add_label(slide, "Verified production runs · 12 Jul 2026", 0.72, 0.93, 3.8)
add_text(slide, "Two missions. Every claim inspectable.", 0.72, 1.22, 11.8, 0.62,
         size=28, font=FONT_SERIF, bold=True)
metric_card(slide, 0.72, 2.03, 2.56, "2", "production briefings", "authenticated Convex workspace")
metric_card(slide, 3.46, 2.03, 2.56, "14 → 8 → 6", "discovered · consumed · published", "full sources beyond snippets")
metric_card(slide, 6.20, 2.03, 2.56, "8", "rejected + preserved", "judgment remains inspectable", accent=OXIDE)
metric_card(slide, 8.94, 2.03, 3.66, "2 exact clips", "timestamp-bounded VideoDB streams", "8.72s + 16.66s", accent=AMBER)
# Receipt card
add_rect(slide, 0.72, 3.42, 7.52, 2.88, fill=GREEN_DARK, radius=True)
add_label(slide, "Evidence receipt", 1.02, 3.71, 1.7, color="9CE1D3")
add_text(slide, "“The result is four recurring failure modes: unregulated growth, missing semantic revision, capacity-driven forgetting, and read-only retrieval.”",
         1.02, 4.08, 6.75, 0.92, size=16, color=WHITE, font=FONT_SERIF, bold=True, line_spacing=1.06)
add_text(slide, "SOURCE  arxiv.org/abs/2605.26252\nLOCATOR Is Agent Memory a Database? · paragraph 13\nSHA-256 f8853e7c…1432385b",
         1.02, 5.21, 6.58, 0.64, size=8.5, color="AFC0BA", font=FONT_MONO, line_spacing=1.18)
add_text(slide, "6 / 6 production claims passed exact evidence validation.", 1.02, 5.94, 5.7, 0.18,
         size=8.5, color="9CE1D3", font=FONT_MONO, bold=True)
# Integrity checklist
add_rect(slide, 8.50, 3.42, 4.10, 2.88, fill=PAPER, line=LINE, radius=True)
add_label(slide, "What the receipts prove", 8.82, 3.71, 2.35)
proofs = [
    "Public app + identity-owned workspace",
    "Live LinkUp discovery + primary fetch",
    "2 Hermes native delegation receipts",
    "6 exact-matched findings; snippets excluded",
    "2 timestamp-bounded VideoDB clips",
]
for i, text in enumerate(proofs):
    y = 4.12 + i * 0.38
    add_rect(slide, 8.84, y, 0.18, 0.18, fill=GREEN, radius=True)
    add_text(slide, "✓", 8.84, y - 0.01, 0.18, 0.18, size=7, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, text, 9.16, y - 0.02, 3.08, 0.23, size=9.5, color=MUTED)
add_text(slide, "HONESTY BOUNDARY · production receipts are real; external users, traffic, and verified model cost are not claimed",
         0.74, 6.61, 11.85, 0.22, size=7.3, color=OXIDE, font=FONT_MONO, bold=True)
add_footer(slide, 6)


# Slide 7 — Close
slide = prs.slides.add_slide(blank)
fill_background(slide, GREEN_DARK)
add_logo(slide, inverse=True)
add_text(slide, "The internet will keep getting bigger.", 0.74, 1.26, 11.4, 0.50,
         size=23, color="AFC0BA", font=FONT_SERIF)
add_text(slide, "Your attention will not.\nCerno decides what deserves it.", 0.74, 1.88, 11.5, 1.48,
         size=36, color=WHITE, font=FONT_SERIF, bold=True, line_spacing=0.94)
# Principles
principles = ["FINITE BRIEFINGS", "VISIBLE TASTE", "CITED EVIDENCE", "NO FEED"]
for i, item in enumerate(principles):
    x = 0.74 + i * 2.96
    add_rect(slide, x, 4.02, 2.64, 0.50, fill="1A3A32", line="315B51", radius=True)
    add_text(slide, item, x + 0.12, 4.15, 2.40, 0.18, size=8.2, color="9CE1D3",
             font=FONT_MONO, bold=True, align=PP_ALIGN.CENTER)
add_line(slide, 0.74, 5.14, 12.60, 5.14, "315B51", 1.0)
add_text(slide, "One mission in. One trusted briefing out.", 0.74, 5.55, 8.8, 0.52,
         size=27, color=WHITE, font=FONT_SERIF, bold=True)
add_text(slide, "cernere · to sift, distinguish, and decide", 0.77, 6.20, 5.4, 0.24,
         size=9, color="AFC0BA", font=FONT_MONO)
add_text(slide, "Q&A", 10.46, 5.56, 2.12, 0.52, size=24, color="9CE1D3", font=FONT_UI,
         bold=True, align=PP_ALIGN.RIGHT)
add_footer(slide, 7, dark=True)


# Presentation metadata
prs.core_properties.title = "Cerno — Your personal research desk, staffed by agents"
prs.core_properties.subject = "GrowthX × Hermes Buildathon — AI as Agency"
prs.core_properties.author = "Cerno"
prs.core_properties.keywords = "Cerno, Hermes, AI as Agency, research agents, finite briefing"
prs.core_properties.comments = "Metrics and evidence are sourced from docs/LIVE-VERIFICATION.md."

out = OUTPUT / "Cerno-Hackathon-Presentation.pptx"
prs.save(out)
print(out)
